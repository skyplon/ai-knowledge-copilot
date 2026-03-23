############################################################
# 🧠 AI KNOWLEDGE COPILOT — V1 (PROFESSIONAL PORTFOLIO VERSION)
############################################################
# OVERVIEW
# ----------------------------------------------------------
# This application is an AI-powered document analysis tool.
#
# It allows users to:
# - Upload multiple file types (PDF, DOCX, PPTX, CSV, TXT)
# - Convert them into a unified text-based knowledge base
# - Ask natural language questions
# - Receive AI-generated answers grounded in those documents
#
# ARCHITECTURE (RAG-LITE)
# ----------------------------------------------------------
# This app uses a simplified Retrieval-Augmented Generation approach:
#
#   Documents → Text Extraction → Context Injection → LLM → Answer
#
# Instead of embeddings/vector DBs, we inject document content
# directly into the LLM prompt (faster, simpler, MVP-friendly).
#
# TECHNOLOGIES USED
# ----------------------------------------------------------
# - Streamlit → UI framework (chat interface + layout)
# - OpenAI GPT-4o-mini → Language model (reasoning engine)
# - PyPDF2 → PDF parsing
# - python-docx → Word document parsing
# - python-pptx → PowerPoint parsing
# - pandas → CSV / structured data handling
#
# DESIGN PRINCIPLES
# ----------------------------------------------------------
# - Simplicity over complexity (no LangChain / vector DB yet)
# - Clean UX (chat-based interaction like ChatGPT)
# - Stateful interaction (chat memory using session_state)
# - Deployment-ready (Streamlit Cloud compatible)
############################################################


############################################################
# 1. IMPORTS — Load Required Libraries
############################################################
# This section imports all dependencies needed for:
# - UI rendering (Streamlit)
# - AI inference (OpenAI)
# - File parsing (PDF, DOCX, PPTX, CSV)
############################################################

import streamlit as st
from openai import OpenAI
import uuid

from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
import pandas as pd


############################################################
# 2. APPLICATION CONFIGURATION
############################################################
# Initializes:
# - Streamlit page settings
# - OpenAI client using secure API key
#
# NOTE:
# The API key must be stored in Streamlit Secrets:
# st.secrets["OPENAI_API_KEY"]
############################################################

st.set_page_config(page_title="AI Knowledge Copilot", layout="wide")

client = OpenAI(api_key=st.secrets["OPENAI_API_KEY"])

############################################################
# 3. ACCESS CONTROL (DEMO PROTECTION)
############################################################
# PURPOSE:
# Restricts access to the app to prevent unauthorized usage
# and protect OpenAI API credits.
#
# HOW IT WORKS:
# - Prompts user for a password
# - Compares against secure value in Streamlit secrets
# - Blocks access if incorrect
#
# NOTE:
# Store password in .streamlit/secrets.toml:
# APP_PASSWORD = "your_password_here"
############################################################

import streamlit as st

# Password input UI
password = st.text_input("🔒 Enter access password to use the app", type="password")

# Stop app if password is incorrect
if password != st.secrets.get("APP_PASSWORD"):
    st.warning("Access restricted. Please enter the correct password.")
    st.stop()


############################################################
# 4. UI STYLING (LIGHTWEIGHT DESIGN SYSTEM)
############################################################
# Applies custom CSS to improve:
# - Spacing
# - Chat bubble appearance
# - Overall readability
#
# This mimics a minimal SaaS / ChatGPT-like UI
############################################################

st.markdown("""
<style>
.block-container { padding-top: 2rem; }
.stChatMessage { border-radius: 12px; padding: 10px; }
</style>
""", unsafe_allow_html=True)


############################################################
# 5. HEADER — Main Application Title
############################################################
# Provides:
# - Product name
# - Short description for users
############################################################

st.title("🧠 AI Knowledge Copilot")
st.markdown("### Analyze documents and interact with your knowledge using AI")


############################################################
# 6. SIDEBAR — WORKSPACE (CONTROL PANEL)
############################################################
# The sidebar acts as the CONTROL PLANE of the app.
#
# Responsibilities:
# - Upload documents
# - Provide suggested questions (UX enhancement)
#
# UX PATTERN:
# Sidebar = Controls
# Main Area = Interaction (chat)
############################################################

with st.sidebar:

    st.header("📂 Workspace")

    st.info("💻 Best experience on desktop. Some mobile-uploaded files may not process correctly (V1 limitation).")
    # File uploader supports multiple file types
    uploaded_files = st.file_uploader(
        "Upload files",
        accept_multiple_files=True,
        type=["pdf", "docx", "pptx", "csv", "txt"]
    )
    st.divider()

    st.subheader("💡 Suggested questions")

    # Session state variable to trigger questions programmatically
    if "trigger_question" not in st.session_state:
        st.session_state.trigger_question = None

    def set_question(q):
        # Stores selected question → injected into chat input
        st.session_state.trigger_question = q

    # Predefined prompts improve usability and onboarding
    st.button("Summarize documents", on_click=set_question, args=("Summarize these documents",))
    st.button("Key insights", on_click=set_question, args=("What are the key insights?",))
    st.button("Risks", on_click=set_question, args=("What are the risks?",))
    st.button("Compare documents", on_click=set_question, args=("Compare these documents",))


############################################################
# 7. DOCUMENT PROCESSING (MULTIMODAL INGESTION)
############################################################
# This function converts different file formats into plain text.
#
# WHY THIS MATTERS:
# LLMs operate on text → all formats must be normalized
#
# Supported formats:
# - PDF → page text extraction
# - DOCX → paragraph extraction
# - PPTX → slide text extraction
# - CSV → structured data → string
# - TXT → raw text
############################################################

def extract_text(file):
    text = ""

    if file.type == "application/pdf":
        reader = PdfReader(file)
        for page in reader.pages:
            text += page.extract_text() or ""

    elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        doc = DocxDocument(file)
        for p in doc.paragraphs:
            text += p.text + "\n"

    elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        prs = Presentation(file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

    elif file.type == "text/csv":
        df = pd.read_csv(file)
        text += df.to_string()

    elif file.type == "text/plain":
        text += file.read().decode("utf-8")

    return text


############################################################
# 8. KNOWLEDGE BASE CONSTRUCTION
############################################################
# Aggregates all uploaded documents into a single context string.
#
# NOTE:
# This replaces vector search in this MVP version.
############################################################

all_text = ""

if uploaded_files:
    with st.spinner("📚 Processing documents..."):
      for file in uploaded_files:
        try:
            all_text += extract_text(file) + "\n\n"
        except Exception as e:
            st.warning(f"⚠️ Could not fully process {file.name}. Try uploading from desktop.")
            continue

    st.success(f"✅ {len(uploaded_files)} files processed successfully")

############################################################
# 9. SESSION STATE (CHAT MEMORY)
############################################################
# Stores conversation history across interactions.
#
# STRUCTURE:
# { id, role, content, feedback }
############################################################

if "messages" not in st.session_state:
    st.session_state.messages = []


############################################################
# 10. CHAT RENDERING (UI)
############################################################
# Displays:
# - Chat messages
# - Feedback buttons (👍 👎)
#
# Uses unique IDs to prevent Streamlit key collisions
############################################################

for msg in st.session_state.messages:

    with st.chat_message(msg["role"]):
        st.markdown(msg["content"])

        if msg["role"] == "assistant":

            msg_id = msg["id"]
            disabled = msg.get("feedback") is not None

            col1, col2 = st.columns(2)

            with col1:
                if st.button("👍", key=f"up_{msg_id}", disabled=disabled):
                    msg["feedback"] = "up"
                    st.rerun()

            with col2:
                if st.button("👎", key=f"down_{msg_id}", disabled=disabled):
                    msg["feedback"] = "down"
                    st.rerun()

            if msg.get("feedback") == "up":
                st.success("Marked as helpful")

            elif msg.get("feedback") == "down":
                st.info("Marked as not helpful")


############################################################
# 11. INPUT HANDLING
############################################################
# Supports:
# - Manual user input
# - Suggested question injection
############################################################

user_input = st.chat_input("Ask anything about your documents...")

if st.session_state.trigger_question:
    user_input = st.session_state.trigger_question
    st.session_state.trigger_question = None


############################################################
# 12. AI RESPONSE GENERATION (LLM CORE)
############################################################
# This is the core intelligence layer.
#
# MODEL:
# - GPT-4o-mini
#
# PROCESS:
# - Inject document context + user question
# - Generate grounded response
############################################################

if user_input:

    if all_text:

        st.session_state.messages.append({
            "id": str(uuid.uuid4()),
            "role": "user",
            "content": user_input
        })

        with st.chat_message("user"):
            st.markdown(user_input)

        with st.chat_message("assistant"):
            with st.spinner("🤖 Thinking..."):

                response = client.chat.completions.create(
                    model="gpt-4o-mini",
                    messages=[
                        {"role": "system", "content": "Answer ONLY using the provided documents."},
                        {"role": "user", "content": f"{all_text[:12000]}\n\nQuestion: {user_input}"}
                    ]
                )

                answer = response.choices[0].message.content
                st.markdown(answer)

        st.session_state.messages.append({
            "id": str(uuid.uuid4()),
            "role": "assistant",
            "content": answer,
            "feedback": None
        })

        st.rerun()

    else:
        st.warning("Please upload documents first.")
